from flask import Flask, request, jsonify, render_template, Response, stream_with_context, redirect, url_for, flash, send_file
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from werkzeug.exceptions import HTTPException
from huggingface_hub import InferenceClient
from dotenv import load_dotenv
import os, json, uuid, io, traceback, re
from datetime import datetime, timezone
from functools import wraps

load_dotenv()  # Charge .env en développement local

app = Flask(__name__, template_folder='../templates', static_folder='../static')

@app.errorhandler(Exception)
def handle_exception(e):
    if isinstance(e, HTTPException):
        return e
    return f"<h1>Internal Server Error</h1><pre>{traceback.format_exc()}</pre>", 500
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'liteai-static-fallback-secret-key-12345')

database_url = os.environ.get('DATABASE_URL', '')
if database_url.startswith("postgres://"):
    database_url = database_url.replace("postgres://", "postgresql://", 1)
app.config['SQLALCHEMY_DATABASE_URI'] = database_url or 'sqlite:///liteai.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['MAX_CONTENT_LENGTH'] = 10 * 1024 * 1024  # 10 Mo max

CORS(app)
db = SQLAlchemy(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'

# ─────────────────────────────────────────────
#  MODÈLES DISPONIBLES
# ─────────────────────────────────────────────
AVAILABLE_MODELS = {
    "qwen3-32b": {
        "id": "Qwen/Qwen3-32B",
        "name": "Qwen3 32B",
        "provider": "novita",
        "description": "Modèle général très puissant, excellent pour la rédaction.",
        "context": "32K tokens",
        "vision": False
    },
    "deepseek-v4-pro": {
        "id": "deepseek-ai/DeepSeek-V4-Pro",
        "name": "DeepSeek V4 Pro",
        "provider": "novita",
        "description": "Raisonnement de pointe et logique complexe.",
        "context": "128K tokens",
        "vision": False
    },
    "deepseek-v4-flash": {
        "id": "deepseek-ai/DeepSeek-V4-Flash",
        "name": "DeepSeek V4 Flash",
        "provider": "novita",
        "description": "Modèle ultra rapide de DeepSeek.",
        "context": "128K tokens",
        "vision": False
    },
    "kimi-k2.6": {
        "id": "moonshotai/Kimi-K2.6",
        "name": "Kimi K2.6",
        "provider": "novita",
        "description": "Excellent modèle conversationnel et analytique.",
        "context": "262K tokens",
        "vision": False
    },
    "glm-5.1": {
        "id": "zai-org/GLM-5.1",
        "name": "GLM 5.1",
        "provider": "novita",
        "description": "Modèle général très polyvalent.",
        "context": "202K tokens",
        "vision": False
    },
    "llama3-70b": {
        "id": "meta-llama/Llama-3.3-70B-Instruct",
        "name": "Llama 3.3 70B",
        "provider": "novita",
        "description": "Excellent pour le code et le raisonnement complexe.",
        "context": "131K tokens",
        "vision": False
    },
    "gpt-oss-120b": {
        "id": "openai/gpt-oss-120b",
        "name": "GPT-OSS 120B",
        "provider": "novita",
        "description": "Open Source State of the Art par OpenAI.",
        "context": "131K tokens",
        "vision": False
    },
    "qwen3-coder-next": {
        "id": "Qwen/Qwen3-Coder-Next",
        "name": "Qwen3 Coder Next",
        "provider": "novita",
        "description": "Le meilleur modèle mondial pour la programmation.",
        "context": "262K tokens",
        "vision": False
    },
    "qwen3-vl": {
        "id": "Qwen/Qwen3-VL-8B-Instruct",
        "name": "Qwen3 VL (Vision)",
        "provider": "novita",
        "description": "Modèle multimodal pour analyser des images.",
        "context": "131K tokens",
        "vision": True
    },
    "deepseek-ocr": {
        "id": "deepseek-ai/DeepSeek-OCR",
        "name": "DeepSeek OCR",
        "provider": "novita",
        "description": "Extraction experte de textes depuis des images.",
        "context": "8K tokens",
        "vision": True
    }
}
DEFAULT_MODEL = "deepseek-v4-pro"

HF_TOKEN = os.environ.get('HF_TOKEN', '')

ALLOWED_EXTENSIONS = {'txt', 'md', 'pdf', 'csv', 'xlsx', 'png', 'jpg', 'jpeg', 'webp'}

# ─────────────────────────────────────────────
#  MODÈLES DE DONNÉES
# ─────────────────────────────────────────────

class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    password_hash = db.Column(db.String(120), nullable=False)
    is_admin = db.Column(db.Boolean, default=False)
    default_model = db.Column(db.String(50), default=DEFAULT_MODEL)
    api_keys = db.relationship('ApiKey', backref='owner', lazy=True)
    sessions = db.relationship('ChatSession', backref='user', lazy=True)

class ApiKey(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    key = db.Column(db.String(50), unique=True, nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    usage_count = db.Column(db.Integer, default=0)
    token_count = db.Column(db.Integer, default=0)
    last_used = db.Column(db.DateTime)
    created_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))

class ChatSession(db.Model):
    id = db.Column(db.String(36), primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    title = db.Column(db.String(100), default="Nouvelle conversation")
    model_key = db.Column(db.String(50), default=DEFAULT_MODEL)
    created_at = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))
    messages = db.relationship('Message', backref='session', lazy=True, cascade="all, delete-orphan")

class Message(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    session_id = db.Column(db.String(36), db.ForeignKey('chat_session.id'), nullable=False)
    role = db.Column(db.String(20), nullable=False)
    content = db.Column(db.Text, nullable=False)
    timestamp = db.Column(db.DateTime, default=lambda: datetime.now(timezone.utc))

# ─────────────────────────────────────────────
#  HELPERS
# ─────────────────────────────────────────────

@login_manager.user_loader
def load_user(user_id):
    return db.session.get(User, int(user_id))

def get_client(model_key):
    model_info = AVAILABLE_MODELS.get(model_key, AVAILABLE_MODELS[DEFAULT_MODEL])
    return InferenceClient(provider=model_info['provider'], api_key=HF_TOKEN), model_info['id']

def require_api_key(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        api_key_str = request.headers.get('X-API-Key') or request.headers.get('x-api-key')
        auth_header = request.headers.get('Authorization', '')
        if auth_header.startswith('Bearer '):
            api_key_str = auth_header.split(' ')[1]
        key_obj = ApiKey.query.filter_by(key=api_key_str).first()
        if not key_obj:
            return jsonify({"error": "Clé API invalide ou manquante"}), 401
        key_obj.usage_count += 1
        key_obj.last_used = datetime.now(timezone.utc)
        db.session.commit()
        request.api_key_obj = key_obj
        request.api_key_owner = key_obj.user_id
        return f(*args, **kwargs)
    return decorated

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_file_content(file_storage):
    """Extraire le texte d'un fichier uploadé."""
    filename = secure_filename(file_storage.filename)
    ext = filename.rsplit('.', 1)[1].lower()
    content = ""
    try:
        if ext in ('txt', 'md'):
            content = file_storage.read().decode('utf-8', errors='ignore')
        elif ext == 'pdf':
            import PyPDF2
            reader = PyPDF2.PdfReader(file_storage)
            content = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext == 'csv':
            import pandas as pd
            df = pd.read_csv(file_storage)
            content = df.to_string()
        elif ext == 'xlsx':
            import pandas as pd
            df = pd.read_excel(file_storage)
            content = df.to_string()
        elif ext in ('png', 'jpg', 'jpeg', 'webp'):
            import base64
            b64_str = base64.b64encode(file_storage.read()).decode('utf-8')
            content = f"data:image/{ext};base64,{b64_str}"
            return filename, content, "image"
    except Exception as e:
        content = f"[Erreur extraction: {e}]"
    return filename, content[:15000], "text"

# ─────────────────────────────────────────────
#  ROUTES UI
# ─────────────────────────────────────────────

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        user = User.query.filter_by(username=username).first()
        if user and check_password_hash(user.password_hash, password):
            login_user(user)
            return redirect(url_for('index'))
        flash('Identifiants invalides')
    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('login'))

@app.route('/emergency-reset')
def emergency_reset():
    try:
        user = User.query.filter_by(username='admin').first()
        new_hash = generate_password_hash('admin123', method='pbkdf2:sha256')
        if user:
            user.password_hash = new_hash
            db.session.commit()
            return "✅ Mot de passe de l'admin réinitialisé avec succès à 'admin123'."
        else:
            admin = User(username='admin', password_hash=new_hash, is_admin=True)
            db.session.add(admin)
            db.session.commit()
            return "✅ Utilisateur 'admin' créé avec succès avec le mot de passe 'admin123'."
    except Exception as e:
        return f"Erreur lors du reset: {str(e)}"

@app.route('/')
@login_required
def index():
    models = {k: {"name": v["name"], "description": v["description"], "context": v["context"]} for k, v in AVAILABLE_MODELS.items()}
    return render_template('index.html', models=models, default_model=current_user.default_model or DEFAULT_MODEL)

@app.route('/docs')
@login_required
def docs():
    return render_template('docs.html', models=AVAILABLE_MODELS)

@app.route('/admin/keys')
@login_required
def admin_keys():
    if not current_user.is_admin:
        return "Accès refusé", 403
    keys = ApiKey.query.order_by(ApiKey.created_at.desc()).all()
    users = User.query.all()
    return render_template('admin_keys.html', keys=keys, users=users, models=AVAILABLE_MODELS)

# ─────────────────────────────────────────────
#  ROUTES API CHAT
# ─────────────────────────────────────────────

@app.route('/api/chat', methods=['POST'])
@require_api_key
def chat():
    data = request.json
    session_id = data.get('session_id')
    model_key = data.get('model', DEFAULT_MODEL)
    if model_key not in AVAILABLE_MODELS:
        model_key = DEFAULT_MODEL
    user_id = request.api_key_owner

    if not session_id:
        session_id = str(uuid.uuid4())
        db.session.add(ChatSession(id=session_id, user_id=user_id, model_key=model_key))
        db.session.commit()

    sess = db.session.get(ChatSession, session_id)
    history = [{"role": m.role, "content": m.content} for m in sess.messages]
    user_content = data.get('message') or (data.get('messages') or [{}])[-1].get('content', '')
    file_type = data.get('file_type')
    file_content = data.get('fileContent')

    if file_type == 'image' and file_content and not AVAILABLE_MODELS.get(model_key, {}).get("vision"):
        return jsonify({"error": "Erreur: Ce modèle ne supporte pas la vision. Choisissez un modèle Vision."}), 400

    db.session.add(Message(session_id=session_id, role='user', content=user_content))
    db.session.commit()

    try:
        ai_client, model_id = get_client(model_key)
        
        messages_payload = history.copy()
        if file_type == 'image' and file_content:
            messages_payload.append({
                "role": "user", 
                "content": [
                    {"type": "text", "text": user_content},
                    {"type": "image_url", "image_url": {"url": file_content}}
                ]
            })
        else:
            messages_payload.append({"role": "user", "content": user_content})

        response = ai_client.chat.completions.create(
            model=model_id,
            messages=messages_payload,
            max_tokens=data.get('max_tokens', 2048),
            temperature=data.get('temperature', 0.7)
        )
        ai_content = response.choices[0].message.content
        db.session.add(Message(session_id=session_id, role='assistant', content=ai_content))
        db.session.commit()
        return jsonify({"session_id": session_id, "response": ai_content, "model": model_key})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/chat/stream', methods=['POST'])
@require_api_key
def chat_stream():
    data = request.json
    session_id = data.get('session_id')
    model_key = data.get('model', DEFAULT_MODEL)
    if model_key not in AVAILABLE_MODELS:
        model_key = DEFAULT_MODEL
    user_id = request.api_key_owner

    if not session_id:
        session_id = str(uuid.uuid4())
        db.session.add(ChatSession(id=session_id, user_id=user_id, model_key=model_key))
        db.session.commit()

    sess = db.session.get(ChatSession, session_id)
    history = [{"role": m.role, "content": m.content} for m in sess.messages]
    user_content = data.get('message') or (data.get('messages') or [{}])[-1].get('content', '')
    file_type = data.get('file_type')
    file_content = data.get('fileContent')

    def generate():
        if file_type == 'image' and file_content and not AVAILABLE_MODELS.get(model_key, {}).get("vision"):
            yield f"data: {json.dumps({'error': 'Erreur: Ce modèle ne supporte pas la vision. Choisissez un modèle Vision.'})}\n\n"
            return

        db.session.add(Message(session_id=session_id, role='user', content=user_content))
        db.session.commit()

        full_ai = ""
        try:
            ai_client, model_id = get_client(model_key)
            
            messages_payload = history.copy()
            if file_type == 'image' and file_content:
                messages_payload.append({
                    "role": "user", 
                    "content": [
                        {"type": "text", "text": user_content},
                        {"type": "image_url", "image_url": {"url": file_content}}
                    ]
                })
            else:
                messages_payload.append({"role": "user", "content": user_content})

            stream = ai_client.chat.completions.create(
                model=model_id,
                messages=messages_payload,
                max_tokens=data.get('max_tokens', 2048),
                temperature=data.get('temperature', 0.7),
                stream=True
            )
            for chunk in stream:
                delta = chunk.choices[0].delta.content
                if delta:
                    full_ai += delta
                    yield f"data: {json.dumps({'content': delta, 'session_id': session_id, 'model': model_key})}\n\n"

            db.session.add(Message(session_id=session_id, role='assistant', content=full_ai))
            if len(sess.messages) <= 2:
                sess.title = user_content[:40] + ("..." if len(user_content) > 40 else "")
            db.session.commit()
            yield "data: [DONE]\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return Response(stream_with_context(generate()), mimetype='text/event-stream')

# ─────────────────────────────────────────────
#  ROUTES COMPAT OpenAI & Anthropic
# ─────────────────────────────────────────────

@app.route('/v1/chat/completions', methods=['POST'])
@require_api_key
def openai_compat_chat():
    data = request.json
    messages = data.get('messages', [])
    model_req = data.get('model', '')
    # Mapper le modèle demandé vers nos clés
    model_key = DEFAULT_MODEL
    for k, v in AVAILABLE_MODELS.items():
        if v['id'] == model_req or k == model_req:
            model_key = k
            break
    try:
        ai_client, model_id = get_client(model_key)
        response = ai_client.chat.completions.create(
            model=model_id,
            messages=messages,
            max_tokens=data.get('max_tokens', 2048),
            temperature=data.get('temperature', 0.7)
        )
        return jsonify({
            "id": f"chatcmpl-{uuid.uuid4().hex[:24]}",
            "object": "chat.completion",
            "model": model_id,
            "choices": [{"index": 0, "message": {"role": "assistant", "content": response.choices[0].message.content}, "finish_reason": "stop"}]
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/v1/messages', methods=['POST'])
@require_api_key
def anthropic_compat_chat():
    data = request.json
    system_prompt = data.get('system', '')
    messages = data.get('messages', [])
    if system_prompt:
        messages.insert(0, {"role": "system", "content": system_prompt})
    try:
        ai_client, model_id = get_client(DEFAULT_MODEL)
        response = ai_client.chat.completions.create(
            model=model_id,
            messages=messages,
            max_tokens=data.get('max_tokens', 2048),
            temperature=data.get('temperature', 0.7)
        )
        ai_content = response.choices[0].message.content
        return jsonify({
            "id": f"msg_{uuid.uuid4().hex[:24]}",
            "type": "message",
            "role": "assistant",
            "model": model_id,
            "content": [{"type": "text", "text": ai_content}],
            "stop_reason": "end_turn",
            "usage": {"input_tokens": 0, "output_tokens": 0}
        })
    except Exception as e:
        return jsonify({"error": {"type": "api_error", "message": str(e)}}), 500

@app.route('/v1/models', methods=['GET'])
@app.route('/api/models', methods=['GET'])
@require_api_key
def list_models():
    data = []
    for key, m in AVAILABLE_MODELS.items():
        data.append({
            "id": m['id'],
            "key": key,
            "object": "model",
            "name": m['name'],
            "description": m['description'],
            "context_window": m['context'],
            "owned_by": "huggingface"
        })
    return jsonify({"object": "list", "data": data})

# ─────────────────────────────────────────────
#  ROUTES UPLOAD / EXPORT
# ─────────────────────────────────────────────

@app.route('/api/upload', methods=['POST'])
@require_api_key
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "Aucun fichier fourni"}), 400
    f = request.files['file']
    if not f or not allowed_file(f.filename):
        return jsonify({"error": f"Format non supporté. Acceptés: {', '.join(ALLOWED_EXTENSIONS)}"}), 400
    filename, content, file_type = extract_file_content(f)
    if not content.strip():
        return jsonify({"error": "Impossible d'extraire le contenu du fichier"}), 422
    return jsonify({"filename": filename, "content": content, "type": file_type, "chars": len(content) if file_type == "text" else 0})

@app.route('/api/export/<fmt>', methods=['POST'])
@require_api_key
def export_file(fmt):
    data = request.json
    content = data.get('content', '')
    filename_base = data.get('filename', 'liteai_export')

    if fmt == 'txt':
        buf = io.BytesIO(content.encode('utf-8'))
        return send_file(buf, mimetype='text/plain', as_attachment=True, download_name=f"{filename_base}.txt")

    elif fmt == 'md':
        buf = io.BytesIO(content.encode('utf-8'))
        return send_file(buf, mimetype='text/markdown', as_attachment=True, download_name=f"{filename_base}.md")

    elif fmt == 'docx':
        from docx import Document
        from docx.shared import Pt
        doc = Document()
        # Parse basic markdown for better design
        for line in content.split('\n'):
            line = line.strip()
            if not line:
                continue
            if line.startswith('# '):
                doc.add_heading(line[2:].replace('**', ''), level=1)
            elif line.startswith('## '):
                doc.add_heading(line[3:].replace('**', ''), level=2)
            elif line.startswith('### '):
                doc.add_heading(line[4:].replace('**', ''), level=3)
            elif line.startswith('- ') or line.startswith('* '):
                doc.add_paragraph(line[2:].replace('**', ''), style='List Bullet')
            else:
                p = doc.add_paragraph()
                parts = line.split('**')
                for i, part in enumerate(parts):
                    run = p.add_run(part)
                    if i % 2 != 0:
                        run.bold = True
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                         as_attachment=True, download_name=f"{filename_base}.docx")

    elif fmt == 'xlsx':
        import pandas as pd
        # Détection de tableaux markdown
        table_lines = [line.strip() for line in content.split('\n') if line.strip().startswith('|') and line.strip().endswith('|')]
        
        if table_lines and len(table_lines) > 2:
            headers = [c.strip() for c in table_lines[0].strip('|').split('|')]
            data = []
            for line in table_lines[2:]:
                row = [c.strip() for c in line.strip('|').split('|')]
                data.append(row)
            # Normaliser la taille des lignes
            for row in data:
                while len(row) < len(headers): row.append("")
                if len(row) > len(headers): del row[len(headers):]
            df = pd.DataFrame(data, columns=headers)
        else:
            # Fallback s'il n'y a pas de tableau explicite
            lines = [{'Contenu': line} for line in content.split('\n')]
            df = pd.DataFrame(lines)

        buf = io.BytesIO()
        with pd.ExcelWriter(buf, engine='openpyxl') as writer:
            df.to_excel(writer, index=False)
            worksheet = writer.sheets['Sheet1']
            for col in worksheet.columns:
                max_length = 0
                column = col[0].column_letter
                for cell in col:
                    try:
                        if len(str(cell.value)) > max_length: max_length = len(cell.value)
                    except: pass
                worksheet.column_dimensions[column].width = min((max_length + 2), 50)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                         as_attachment=True, download_name=f"{filename_base}.xlsx")

    elif fmt == 'pptx':
        from pptx import Presentation
        prs = Presentation()
        # Séparer par "---" (format standard de slide markdown) ou par titres "##"
        sections = re.split(r'\n---\n', content)
        if len(sections) == 1:
            parts = re.split(r'\n##\s+', '\n' + content)
            sections = ["## " + p for p in parts[1:]] if len(parts) > 1 else [content]
        
        for section in sections:
            lines = [l.strip() for l in section.split('\n') if l.strip()]
            if not lines: continue
            
            title = "Diapositive"
            bullets = []
            
            if lines[0].startswith('# '):
                title = lines[0][2:].replace('**', '')
                lines = lines[1:]
            elif lines[0].startswith('## '):
                title = lines[0][3:].replace('**', '')
                lines = lines[1:]
                
            for line in lines:
                line = line.replace('**', '')
                if line.startswith('- ') or line.startswith('* '):
                    bullets.append(line[2:])
                elif line:
                    bullets.append(line)
                    
            # Layout Titre + Contenu
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = title
            tf = body_shape.text_frame
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(buf, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True, download_name=f"{filename_base}.pptx")

    elif fmt == 'pdf':
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        for line in content.split('\n'):
            if line.startswith('# '):
                story.append(Paragraph(line[2:].replace('**', ''), styles['Heading1']))
            elif line.startswith('## '):
                story.append(Paragraph(line[3:].replace('**', ''), styles['Heading2']))
            elif line:
                story.append(Paragraph(line.replace('**', ''), styles['Normal']))
            story.append(Spacer(1, 5))
        doc.build(story)
        buf.seek(0)
        return send_file(buf, mimetype='application/pdf', as_attachment=True, download_name=f"{filename_base}.pdf")

    return jsonify({"error": f"Format '{fmt}' non supporté. Acceptés: txt, md, docx, xlsx, pptx, pdf"}), 400

# ─────────────────────────────────────────────
#  ROUTES SESSIONS & ADMIN
# ─────────────────────────────────────────────

@app.route('/api/sessions', methods=['GET'])
@login_required
def list_sessions():
    sessions = ChatSession.query.filter_by(user_id=current_user.id).order_by(ChatSession.created_at.desc()).all()
    return jsonify([{
        "id": s.id, "title": s.title,
        "model": s.model_key,
        "model_name": AVAILABLE_MODELS.get(s.model_key, {}).get('name', s.model_key),
        "date": s.created_at.isoformat()
    } for s in sessions])

@app.route('/api/history/<session_id>', methods=['GET'])
@login_required
def get_session_history(session_id):
    sess = db.session.get(ChatSession, session_id)
    if not sess: return jsonify({"error": "Introuvable"}), 404
    if sess.user_id != current_user.id: return jsonify({"error": "Interdit"}), 403
    return jsonify({
        "id": sess.id, "title": sess.title,
        "model": sess.model_key,
        "messages": [{"role": m.role, "content": m.content} for m in sess.messages]
    })

@app.route('/api/sessions/<session_id>', methods=['DELETE'])
@login_required
def delete_session(session_id):
    sess = db.session.get(ChatSession, session_id)
    if not sess or sess.user_id != current_user.id: return jsonify({"error": "Interdit"}), 403
    db.session.delete(sess)
    db.session.commit()
    return jsonify({"success": True})

@app.route('/api/admin/keys', methods=['POST'])
@login_required
def create_key():
    if not current_user.is_admin: return jsonify({"error": "Interdit"}), 403
    data = request.json or {}
    new_key = f"sk-liteai-{uuid.uuid4().hex[:16]}"
    target_user_id = data.get('user_id', current_user.id)
    db.session.add(ApiKey(key=new_key, user_id=target_user_id))
    db.session.commit()
    return jsonify({"success": True, "key": new_key})

@app.route('/api/admin/keys/<key>', methods=['DELETE'])
@login_required
def delete_key(key):
    if not current_user.is_admin: return jsonify({"error": "Interdit"}), 403
    key_obj = ApiKey.query.filter_by(key=key).first()
    if key_obj:
        db.session.delete(key_obj)
        db.session.commit()
        return jsonify({"success": True})
    return jsonify({"error": "Clé introuvable"}), 404

@app.route('/api/admin/users', methods=['POST'])
@login_required
def create_user():
    if not current_user.is_admin: return jsonify({"error": "Interdit"}), 403
    data = request.json or {}
    username = data.get('username', '').strip()
    password = data.get('password', '')
    if not username or not password:
        return jsonify({"error": "Champs manquants"}), 400
    if User.query.filter_by(username=username).first():
        return jsonify({"error": "Utilisateur existe déjà"}), 400
    new_user = User(username=username, password_hash=generate_password_hash(password, method='pbkdf2:sha256'), is_admin=data.get('is_admin', False))
    db.session.add(new_user)
    db.session.commit()
    return jsonify({"success": True, "id": new_user.id})

@app.route('/api/admin/stats', methods=['GET'])
@login_required
def get_stats():
    if not current_user.is_admin: return jsonify({"error": "Interdit"}), 403
    return jsonify({
        "total_users": User.query.count(),
        "total_keys": ApiKey.query.count(),
        "total_sessions": ChatSession.query.count(),
        "total_messages": Message.query.count(),
        "total_requests": db.session.query(db.func.sum(ApiKey.usage_count)).scalar() or 0,
    })

# ─────────────────────────────────────────────
#  INIT
# ─────────────────────────────────────────────

with app.app_context():
    try:
        db.create_all()
        if not User.query.first():
            admin = User(username='admin', password_hash=generate_password_hash('admin123', method='pbkdf2:sha256'), is_admin=True)
            db.session.add(admin)
            db.session.commit()
            db.session.add(ApiKey(key='sk-liteai-12345', user_id=admin.id))
            db.session.commit()
            print("✅ Admin créé: admin / admin123")
    except Exception as e:
        print(f"[DB Init Error] {e}")

if __name__ == '__main__':
    app.run(debug=True, port=5000)
